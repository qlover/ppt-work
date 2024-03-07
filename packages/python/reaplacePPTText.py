from pptx import Presentation
import os

# # 加载现有的PPT文件
# ppt = Presentation("D:\\qrj\\download\\2.29\\01 详情展示\\企业宣传- (7).pptx")

# # 遍历每个幻灯片
# for slide in ppt.slides:
#     # 遍历每个形状
#     for shape in slide.shapes:
#         # 检查形状是否是文本框
#         if shape.has_text_frame:
#             # 遍历文本框中的每个段落
#             for paragraph in shape.text_frame.paragraphs:
#                 # 遍历段落中的每个文本块
#                 for run in paragraph.runs:
#                     # 替换文本
#                     run.text = run.text.replace("唐峰", "晓橙")

# # 保存修改后的PPT文件
# ppt.save("D:\\qrj\\download\\2.29\\01 详情展示\\企业宣传- (7).pptx")



def replace_pptx(directory):
    
    # 遍历目录下的所有文件
    for filename in os.listdir(directory):
        # 检查文件是否是PPT文件
        if filename.endswith(".pptx"):
            # 加载PPT文件
            ppt = Presentation(os.path.join(directory, filename))
            
            # 遍历每个幻灯片
            for slide in ppt.slides:
                # 遍历每个形状
                for shape in slide.shapes:
                    # 检查形状是否是文本框
                    if shape.has_text_frame:
                        # 遍历文本框中的每个段落
                        for paragraph in shape.text_frame.paragraphs:
                            # 遍历段落中的每个文本块
                            for run in paragraph.runs:
                                # 替换文本
                                run.text = run.text.replace("唐峰", "晓橙")
                                run.text = run.text.replace("创妹", "晓橙")
                                run.text = run.text.replace("tangfeng", "xiaocheng")
            
            # 保存修改后的PPT文件
            ppt.save(os.path.join(directory, filename))


# 定义目录路径
# replace_pptx('D:\\qrj\\download\\2.29\\22图标系列')
# replace_pptx('D:\\qrj\\download\\2.29\\23图表风格')
# replace_pptx('D:\\qrj\\download\\2.29\\24相册纪念')
# replace_pptx('D:\\qrj\\download\\2.29\\25星空风格')
# replace_pptx('D:\\qrj\\download\\2.29\\7教师课件')
# replace_pptx('D:\\qrj\\download\\2.29\\8简历求职')
# replace_pptx('D:\\qrj\\download\\2.29\\9清新文艺')
# replace_pptx('D:\\qrj\\download\\3.1\\01 述职报告')

# replace_pptx('D:\\qrj\\download\\3.1\\05 年中总结')