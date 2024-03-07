import os
from pptx import Presentation


def replace_text_in_pptx(pptx_file, new_text):
    try:
        prs = Presentation(pptx_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = run.text.replace("唐峰", new_text)
                            run.text = run.text.replace("创妹", new_text)
                            run.text = run.text.replace("tangfeng", new_text)
        prs.save(pptx_file)
        print(f"Successfully processed: {pptx_file}")
    except Exception as e:
        print(f"Error processing {pptx_file}: {e}")


def replace_text_in_directory(directory, new_text):
    for filename in os.listdir(directory):
        if filename.endswith(".pptx"):
            file_path = os.path.join(directory, filename)
            replace_text_in_pptx(file_path, new_text)


# 指定目录
directory_path = "D:\\qrj\\download\\3.1\\00 详情展示部分\\02 年中总结更新"
new_text = "晓橙"

replace_text_in_directory(directory_path, new_text)
